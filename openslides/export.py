"""
Export Engine
HTML slides -> PDF (Chrome headless), PNG (Playwright), PPTX.
"""
from __future__ import annotations

import asyncio
import shutil
import subprocess
import tempfile
from pathlib import Path

from PyPDF2 import PdfReader, PdfWriter
from PyPDF2.generic import (
    ArrayObject, DictionaryObject, NameObject, NumberObject, TextStringObject,
)


# CSS injected into every slide before PDF rendering to fix Chrome headless
# rendering artifacts: box-shadow renders as gray blocks, radial-gradient
# on pseudo-elements renders as solid rectangles.
_PRINT_FIX_CSS = """
<style>
@media print {
  * { box-shadow: none !important; -webkit-box-shadow: none !important; }
  *::before, *::after {
    background-image: none !important;
    box-shadow: none !important;
  }
}
</style>
"""


def _find_chrome() -> str | None:
    """Find Chrome binary on the system, or None if not available."""
    for name in ("google-chrome-stable", "google-chrome", "chromium-browser", "chromium"):
        path = shutil.which(name)
        if path:
            return path
    return None


def _ensure_playwright_browsers():
    """Install Playwright Chromium + system deps if not already present."""
    cache = Path.home() / ".cache" / "ms-playwright"
    if not any(cache.glob("chromium*")) if cache.exists() else True:
        subprocess.run(
            ["playwright", "install", "--with-deps", "chromium"],
            check=True, capture_output=True,
        )


def _render_slide_pdf_playwright(
    html: str, pdf_path: Path, width: int, height: int
) -> Path:
    """Render a single slide HTML to PDF using Playwright (sync API)."""
    from playwright.sync_api import sync_playwright

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": width, "height": height})
        page.set_content(html, wait_until="networkidle")
        page.evaluate("() => document.fonts.ready")
        page.pdf(
            path=str(pdf_path),
            width=f"{width}px",
            height=f"{height}px",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )
        browser.close()
    return pdf_path


def _inject_print_css(html: str) -> str:
    """Inject @media print CSS fixes before </head> in slide HTML."""
    if "</head>" in html:
        return html.replace("</head>", _PRINT_FIX_CSS + "</head>", 1)
    if "<body" in html:
        return html.replace("<body", _PRINT_FIX_CSS + "<body", 1)
    return _PRINT_FIX_CSS + html


def _inject_page_size(html: str, width: int, height: int) -> str:
    """Inject @page size rule if not already present."""
    if "@page" in html:
        return html
    page_css = f"<style>@page {{ size: {width}px {height}px; margin: 0; }}</style>"
    if "</head>" in html:
        return html.replace("</head>", page_css + "</head>", 1)
    return page_css + html


def _render_slide_pdf(
    chrome: str, html_path: Path, pdf_path: Path, timeout: int = 30
) -> Path:
    """Render a single HTML file to PDF using Chrome headless."""
    cmd = [
        chrome,
        "--headless",
        "--no-sandbox",
        "--disable-gpu",
        "--disable-software-rasterizer",
        "--disable-dev-shm-usage",
        "--font-render-hinting=none",
        f"--print-to-pdf={pdf_path}",
        "--print-to-pdf-no-header",
        "--no-margins",
        str(html_path),
    ]
    result = subprocess.run(
        cmd, capture_output=True, text=True, timeout=timeout
    )
    if not pdf_path.exists():
        raise RuntimeError(
            f"Chrome headless PDF failed for {html_path}.\n"
            f"stderr: {result.stderr}\nstdout: {result.stdout}"
        )
    return pdf_path


def export_pdf_sync(
    slides: list[str],
    output_path: str | Path,
    links: list[dict] | None = None,
    width: int = 1920,
    height: int = 1080,
) -> Path:
    """
    Export slides as a single merged PDF using Chrome headless.

    Each slide HTML is written to a temp file, rendered to PDF via
    google-chrome --headless --print-to-pdf, then merged with PyPDF2.

    Args:
        slides: list of HTML strings (full documents)
        output_path: destination PDF path
        links: optional list of {"page": int, "x1", "y1", "x2", "y2", "url"}
        width: slide width in pixels (default 1920)
        height: slide height in pixels (default 1080)
    """
    output_path = Path(output_path).expanduser()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    chrome = _find_chrome()
    use_playwright = chrome is None
    if use_playwright:
        _ensure_playwright_browsers()

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        pdf_paths = []

        for i, html in enumerate(slides):
            html = _inject_page_size(html, width, height)
            html = _inject_print_css(html)

            pdf_path = tmpdir / f"slide-{i:02d}.pdf"

            if use_playwright:
                _render_slide_pdf_playwright(html, pdf_path, width, height)
            else:
                html_path = tmpdir / f"slide-{i:02d}.html"
                html_path.write_text(html, encoding="utf-8")
                _render_slide_pdf(chrome, html_path, pdf_path)

            pdf_paths.append(pdf_path)

        # Merge all single-page PDFs
        writer = PdfWriter()
        for pdf_path in pdf_paths:
            for pg in PdfReader(str(pdf_path)).pages:
                writer.add_page(pg)

        # Add link annotations
        if links:
            for link in links:
                _add_pdf_link(
                    writer,
                    page_index=link["page"],
                    x1=link["x1"],
                    y1=link["y1"],
                    x2=link["x2"],
                    y2=link["y2"],
                    url=link["url"],
                )

        with open(output_path, "wb") as f:
            writer.write(f)

    return output_path


class ExportEngine:
    """
    Async export engine for PNG (uses Playwright).
    PDF export uses Chrome headless directly (see export_pdf_sync).

    Usage:
        async with ExportEngine() as engine:
            pngs = await engine.export_png(slides, "./pngs/")
    """

    def __init__(self, width: int = 1920, height: int = 1080):
        self.width = width
        self.height = height
        self._pw = None
        self._browser = None

    async def _browser_instance(self):
        if self._browser is None:
            from playwright.async_api import async_playwright
            self._pw = await async_playwright().start()
            launch_kwargs = {"headless": True}
            chrome_path = _find_chrome()
            if chrome_path:
                launch_kwargs["executable_path"] = chrome_path
            self._browser = await self._pw.chromium.launch(**launch_kwargs)
        return self._browser

    async def export_png(
        self,
        slides: list[str],
        output_dir: str | Path,
        prefix: str = "slide",
    ) -> list[Path]:
        """Export each slide as PNG."""
        output_dir = Path(output_dir).expanduser()
        output_dir.mkdir(parents=True, exist_ok=True)

        browser = await self._browser_instance()
        paths = []

        for i, html in enumerate(slides, 1):
            page = await browser.new_page(
                viewport={"width": self.width, "height": self.height}
            )
            await page.set_content(html, wait_until="networkidle")
            await page.evaluate("() => document.fonts.ready")

            path = output_dir / f"{prefix}-{i:02d}.png"
            await page.screenshot(path=str(path), full_page=False, type="png")
            await page.close()
            paths.append(path)

        return paths

    async def close(self):
        if self._browser:
            await self._browser.close()
            self._browser = None
        if self._pw:
            await self._pw.stop()
            self._pw = None

    async def __aenter__(self):
        return self

    async def __aexit__(self, *_):
        await self.close()


def export_png_sync(slides: list[str], output_dir: str | Path) -> list[Path]:
    """Synchronous PNG export."""
    return asyncio.run(_async_export_png(slides, output_dir))


async def _async_export_png(slides, output_dir):
    async with ExportEngine() as engine:
        return await engine.export_png(slides, output_dir)


def export_pptx(slides: list[str], output_path: str | Path) -> Path:
    """
    Export slides as PPTX.

    Creates a simplified PPTX from HTML slides. Each slide becomes an image
    embedded in a PowerPoint slide (preserves exact visual layout).
    """
    from pptx import Presentation
    from pptx.util import Inches, Emu

    output_path = Path(output_path).expanduser()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # First render to PNGs
    with tempfile.TemporaryDirectory() as tmpdir:
        pngs = export_png_sync(slides, tmpdir)

        prs = Presentation()
        # Set slide dimensions to 16:9
        prs.slide_width = Emu(12192000)   # 10 inches * 914400
        prs.slide_height = Emu(6858000)   # 7.5 inches * 914400

        blank_layout = prs.slide_layouts[6]  # blank layout

        for png_path in pngs:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(png_path),
                left=Emu(0),
                top=Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

        prs.save(str(output_path))

    return output_path


def _add_pdf_link(writer: PdfWriter, page_index: int, x1: float, y1: float, x2: float, y2: float, url: str):
    """Add a clickable link annotation to a PDF page."""
    if page_index >= len(writer.pages):
        return
    page = writer.pages[page_index]
    link = DictionaryObject({
        NameObject("/Type"): NameObject("/Annot"),
        NameObject("/Subtype"): NameObject("/Link"),
        NameObject("/Rect"): ArrayObject([
            NumberObject(int(x1)), NumberObject(int(y1)),
            NumberObject(int(x2)), NumberObject(int(y2)),
        ]),
        NameObject("/Border"): ArrayObject([
            NumberObject(0), NumberObject(0), NumberObject(0),
        ]),
        NameObject("/A"): DictionaryObject({
            NameObject("/S"): NameObject("/URI"),
            NameObject("/URI"): TextStringObject(url),
        }),
    })
    if "/Annots" not in page:
        page[NameObject("/Annots")] = ArrayObject()
    page[NameObject("/Annots")].append(link)
